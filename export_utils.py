"""
export_utils.py
---------------
Generates downloadable exports from HMIS query results and chart images.

  generate_excel(title, rows, upload_date) -> (BytesIO, base_filename)
  generate_ppt(title, chart_buf, upload_date)  -> (BytesIO, base_filename)
"""

import io
import logging
import re
from datetime import datetime

import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.utils import get_column_letter

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.enum.chart import (
    XL_CHART_TYPE, XL_LEGEND_POSITION,
    XL_LABEL_POSITION, XL_MARKER_STYLE,
)

logger = logging.getLogger(__name__)

# ── Colour palette (matches chart.py) ────────────────────────────────────────
_NAVY_HEX  = "003D6B"
_NAVY_RGB  = RGBColor(0x00, 0x3D, 0x6B)
_ALT_HEX   = "F0F4F8"   # light blue-grey alternating row tint
_SRC_HEX   = "E8EEF4"   # source-line background
_RULE_RGB  = RGBColor(0xDD, 0xDD, 0xDD)
_GRAY_RGB  = RGBColor(0xAA, 0xAA, 0xAA)
_TEXT_HEX  = "1A1A1A"

_FOOTER_BASE = "Source: Bihar HMIS Data  |  HMIS Analytics Bot"

# McKinsey palette as RGBColor objects (mirrors chart.py _PALETTE)
_PALETTE_RGB = [
    RGBColor(0x00, 0x3D, 0x6B),  # navy
    RGBColor(0x2E, 0x86, 0xC1),  # steel
    RGBColor(0x14, 0x8A, 0x72),  # teal
    RGBColor(0xD4, 0x86, 0x0B),  # amber
    RGBColor(0x7B, 0x3F, 0xA0),  # purple
    RGBColor(0x15, 0x65, 0xA8),  # blue
]

# Chart type strings that can be rendered as native editable pptx charts
_NATIVE_TYPES: dict = {
    "bar":                    XL_CHART_TYPE.COLUMN_CLUSTERED,
    "horizontal_bar":         XL_CHART_TYPE.BAR_CLUSTERED,
    "grouped_bar":            XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line":                   XL_CHART_TYPE.LINE_MARKERS,
    "area":                   XL_CHART_TYPE.AREA,
    "stacked_area":           XL_CHART_TYPE.AREA_STACKED,
    "stacked_bar":            XL_CHART_TYPE.COLUMN_STACKED,
    "stacked_column":         XL_CHART_TYPE.COLUMN_STACKED,
    "stacked_horizontal_bar": XL_CHART_TYPE.BAR_STACKED,
    "pie":                    XL_CHART_TYPE.PIE,
    "donut":                  XL_CHART_TYPE.DOUGHNUT,
    "radar":                  XL_CHART_TYPE.RADAR_MARKERS,
    "scatter":                XL_CHART_TYPE.XY_SCATTER,
    "connected_scatter":      XL_CHART_TYPE.XY_SCATTER_LINES,
    "bubble":                 XL_CHART_TYPE.BUBBLE,
    "slope":                  XL_CHART_TYPE.LINE_MARKERS,
}

_LINE_TYPES = {"line", "slope", "connected_scatter"}
_AREA_TYPES = {"area", "stacked_area"}
_BAR_TYPES  = {
    "bar", "horizontal_bar", "grouped_bar",
    "stacked_bar", "stacked_column", "stacked_horizontal_bar",
}


# ── Shared helpers ────────────────────────────────────────────────────────────

def _safe_filename(title: str) -> str:
    """Convert a chart title to a safe ASCII filename base (no extension)."""
    safe = re.sub(r"[^\w\s\-]", "", title).strip()
    safe = re.sub(r"\s+", "_", safe)
    return f"HMIS_{safe[:44]}_{datetime.now().strftime('%Y%m%d')}"


def _clean_col(name: str) -> str:
    """Turn a BigQuery column name into a readable header label."""
    name = re.sub(r"^[_m\d]+", "", name)   # strip leading _9_1_3_ / m4_ prefixes
    name = name.replace("_", " ").strip()
    return name.title() if name else name


def _footer_text(upload_date: str | None) -> str:
    parts = [_FOOTER_BASE]
    if upload_date:
        parts.append(f"Data upload date: {upload_date}")
    parts.append(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    return "  |  ".join(parts)


# ── Excel ─────────────────────────────────────────────────────────────────────

def generate_excel(
    title: str,
    rows: list[dict],
    upload_date: str | None,
) -> tuple[io.BytesIO, str]:
    """
    Build a clean McKinsey-style .xlsx from BigQuery result rows.
    Returns (BytesIO of the file, base filename without extension).
    """
    wb  = openpyxl.Workbook()
    ws  = wb.active
    ws.title = "HMIS Data"

    columns = list(rows[0].keys()) if rows else []
    n_cols  = max(len(columns), 1)

    navy_fill  = PatternFill("solid", fgColor=_NAVY_HEX)
    alt_fill   = PatternFill("solid", fgColor=_ALT_HEX)
    white_fill = PatternFill("solid", fgColor="FFFFFF")
    src_fill   = PatternFill("solid", fgColor=_SRC_HEX)

    def _merge(row_n: int) -> None:
        if n_cols > 1:
            ws.merge_cells(start_row=row_n, start_column=1,
                           end_row=row_n,   end_column=n_cols)

    # Row 1 — Title
    ws.append([title] + [""] * (n_cols - 1))
    _merge(1)
    c = ws.cell(1, 1)
    c.font      = Font(name="Trebuchet MS", bold=True, color="FFFFFF", size=14)
    c.fill      = navy_fill
    c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    ws.row_dimensions[1].height = 30

    # Row 2 — Source line
    ws.append([_FOOTER_BASE] + [""] * (n_cols - 1))
    _merge(2)
    c = ws.cell(2, 1)
    c.font      = Font(name="Trebuchet MS", italic=True, color="555555", size=9)
    c.fill      = src_fill
    c.alignment = Alignment(horizontal="left", vertical="center")
    ws.row_dimensions[2].height = 15

    # Row 3 — Spacer
    ws.append([""] * n_cols)
    ws.row_dimensions[3].height = 6

    # Row 4 — Column headers
    ws.append([_clean_col(col) for col in columns])
    ws.row_dimensions[4].height = 22
    for ci in range(1, n_cols + 1):
        cell = ws.cell(4, ci)
        cell.font      = Font(name="Trebuchet MS", bold=True, color="FFFFFF", size=10)
        cell.fill      = navy_fill
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)

    # Rows 5+ — Data
    data_font = Font(name="Trebuchet MS", color=_TEXT_HEX, size=10)
    for ri, row_dict in enumerate(rows, start=5):
        ws.append([row_dict.get(col) for col in columns])
        ws.row_dimensions[ri].height = 18
        fill = alt_fill if ri % 2 == 1 else white_fill
        for ci in range(1, n_cols + 1):
            cell = ws.cell(ri, ci)
            cell.font      = data_font
            cell.fill      = fill
            cell.alignment = Alignment(horizontal="left", vertical="center")
            if isinstance(cell.value, float):
                cell.number_format = "#,##0.0"
            elif isinstance(cell.value, int):
                cell.number_format = "#,##0"

    # Footer rows (spacer + footer text)
    footer_row = len(rows) + 6
    ws.append([""] * n_cols)
    ws.append([_footer_text(upload_date)] + [""] * (n_cols - 1))
    _merge(footer_row + 1)
    c = ws.cell(footer_row + 1, 1)
    c.font      = Font(name="Trebuchet MS", italic=True, color="AAAAAA", size=8)
    c.alignment = Alignment(horizontal="left", vertical="center")

    # Auto column widths
    for ci, col_name in enumerate(columns, 1):
        header_len   = len(_clean_col(col_name))
        max_data_len = max((len(str(r.get(col_name) or "")) for r in rows), default=0)
        ws.column_dimensions[get_column_letter(ci)].width = min(max(header_len, max_data_len) + 3, 42)

    ws.freeze_panes = "A5"

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf, _safe_filename(title)


# ── PowerPoint ────────────────────────────────────────────────────────────────

def _to_num(v) -> float | None:
    """Coerce any value to float, returning None on failure."""
    try:
        return float(v)
    except (TypeError, ValueError):
        return None


def _make_category_chart_data(
    spec: dict,
    rows: list[dict],
) -> "CategoryChartData | None":
    """Build CategoryChartData from spec + rows (wide or long/pivot format)."""
    x_col      = spec.get("x_column")
    y_cols     = spec.get("y_columns") or []
    y_labels   = spec.get("y_labels") or y_cols
    series_col = spec.get("series_column")

    if not x_col or not y_cols or not rows:
        return None

    chart_data = CategoryChartData()

    if series_col:
        # Long-format: pivot rows on series_col → one series per unique series value
        y_col = y_cols[0]
        cats = list(dict.fromkeys(str(r.get(x_col, "")) for r in rows))
        series_vals = list(dict.fromkeys(str(r.get(series_col, "")) for r in rows))
        chart_data.categories = cats
        lookup = {
            (str(r.get(series_col, "")), str(r.get(x_col, ""))): _to_num(r.get(y_col))
            for r in rows
        }
        for sv in series_vals:
            vals = tuple(lookup.get((sv, c)) for c in cats)
            chart_data.add_series(sv, vals)
    else:
        # Wide-format: one series per y_column
        chart_data.categories = [str(r.get(x_col, "")) for r in rows]
        for y_col, label in zip(y_cols, y_labels):
            vals = tuple(_to_num(r.get(y_col)) for r in rows)
            chart_data.add_series(str(label), vals)

    return chart_data


def _make_xy_chart_data(
    spec: dict,
    rows: list[dict],
    bubble: bool = False,
):
    """Build XyChartData or BubbleChartData for scatter / bubble charts."""
    x_col    = spec.get("x_column")
    y_cols   = spec.get("y_columns") or []
    y_labels = spec.get("y_labels") or y_cols
    size_col = spec.get("size_column")

    if not x_col or not y_cols or not rows:
        return None

    if bubble and size_col:
        chart_data = BubbleChartData()
        for y_col, label in zip(y_cols, y_labels):
            series = chart_data.add_series(str(label))
            for r in rows:
                xv = _to_num(r.get(x_col))
                yv = _to_num(r.get(y_col))
                sv = _to_num(r.get(size_col)) or 1.0
                if xv is not None and yv is not None:
                    series.add_data_point(xv, yv, sv)
    else:
        chart_data = XyChartData()
        for y_col, label in zip(y_cols, y_labels):
            series = chart_data.add_series(str(label))
            for r in rows:
                xv = _to_num(r.get(x_col))
                yv = _to_num(r.get(y_col))
                if xv is not None and yv is not None:
                    series.add_data_point(xv, yv)

    return chart_data


def _style_chart(chart, chart_type_str: str) -> None:
    """
    Replicate chart.py's McKinsey visual style on a native pptx chart object:
      - Trebuchet MS, 10 pt, #555555 on all axis labels and legend
      - McKinsey colour palette on series fills / lines
      - White chart area and plot area backgrounds
      - Very light gridlines (#F0F0F0, 0.8pt) on value axis only
      - #DDDDDD spine / axis line colour
      - Circle markers (white fill, coloured border, size 7) for line charts
      - Data labels: INSIDE_END/white/bold for bars, CENTER/white for stacked,
        ABOVE/dark for line charts
    """
    _TEXT_RGB   = RGBColor(0x1A, 0x1A, 0x1A)
    _TICK_RGB   = RGBColor(0x55, 0x55, 0x55)
    _SPINE_RGB  = RGBColor(0xDD, 0xDD, 0xDD)
    _GRID_RGB   = RGBColor(0xF0, 0xF0, 0xF0)
    _WHITE_RGB  = RGBColor(0xFF, 0xFF, 0xFF)
    _BORDER_RGB = RGBColor(0xCC, 0xCC, 0xCC)

    is_line    = chart_type_str in _LINE_TYPES
    is_area    = chart_type_str in _AREA_TYPES
    is_bar     = chart_type_str in {"bar", "horizontal_bar", "grouped_bar"}
    is_stacked = chart_type_str in {"stacked_bar", "stacked_column", "stacked_horizontal_bar"}
    is_pie     = chart_type_str in ("pie", "donut")
    is_scatter = chart_type_str in ("scatter", "connected_scatter", "bubble")

    # ── 1. White backgrounds + thin border (matches figure border in chart.py) ─
    for area_attr in ("chart_area", "plot_area"):
        try:
            area = getattr(chart, area_attr)
            area.fill.solid()
            area.fill.fore_color.rgb = _WHITE_RGB
        except Exception:
            pass
    try:
        chart.chart_area.format.line.color.rgb = _BORDER_RGB
        chart.chart_area.format.line.width = Pt(0.9)
    except Exception:
        pass

    # ── 2. Hide built-in chart title (slide has its own textbox) ─────────────
    try:
        chart.has_title = False
    except Exception:
        pass

    # ── 3. Legend: bottom, Trebuchet MS 10pt, no frame ───────────────────────
    try:
        n_series = len(list(chart.series))
        chart.has_legend = n_series > 1
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.name = "Trebuchet MS"
            chart.legend.font.size = Pt(10)
            chart.legend.font.color.rgb = _TEXT_RGB
    except Exception:
        pass

    # ── 4. Series colours, line widths, and markers ───────────────────────────
    try:
        for i, series in enumerate(chart.series):
            colour = _PALETTE_RGB[i % len(_PALETTE_RGB)]

            if is_line or is_scatter:
                # 2.5pt line + white-fill circle marker with coloured border
                series.format.line.color.rgb = colour
                series.format.line.width = Pt(2.5)
                series.marker.style = XL_MARKER_STYLE.CIRCLE
                series.marker.size = 7
                series.marker.format.fill.solid()
                series.marker.format.fill.fore_color.rgb = _WHITE_RGB
                series.marker.format.line.color.rgb = colour
                series.marker.format.line.width = Pt(2.2)
            elif not is_pie:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = colour
                if is_area:
                    # Area charts also carry a thin line on top of the fill
                    series.format.line.color.rgb = colour
                    series.format.line.width = Pt(1.5)
    except Exception:
        pass

    # ── 5. Data labels ────────────────────────────────────────────────────────
    try:
        plot = chart.plots[0]
        plot.has_data_labels = True
        dls = plot.data_labels
        dls.font.name = "Trebuchet MS"
        dls.font.bold = True
        if is_bar:
            # Inside bar end (white), PowerPoint auto-flips outside when bar is too short
            dls.number_format = "0.0"
            dls.font.size = Pt(10)
            dls.font.color.rgb = _WHITE_RGB
            dls.position = XL_LABEL_POSITION.INSIDE_END
        elif is_stacked:
            # Segment centre (white), skip tiny segments visually in PPT
            dls.number_format = "0"
            dls.font.size = Pt(9)
            dls.font.color.rgb = _WHITE_RGB
            dls.position = XL_LABEL_POSITION.CENTER
        elif is_line:
            # Above each marker point (dark text) — matches chart.py annotate()
            dls.number_format = "0.0"
            dls.font.size = Pt(9)
            dls.font.color.rgb = _TEXT_RGB
            dls.position = XL_LABEL_POSITION.ABOVE
        # pie / donut / scatter / radar: no programmatic labels (too cluttered)
    except Exception:
        pass

    # ── 6. Value axis: light gridlines, Trebuchet MS 10pt, #DDDDDD spine ─────
    try:
        va = chart.value_axis
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = _GRID_RGB
        va.major_gridlines.format.line.width = Pt(0.8)
        va.has_minor_gridlines = False
        va.format.line.color.rgb = _SPINE_RGB
        va.tick_labels.font.name = "Trebuchet MS"
        va.tick_labels.font.size = Pt(10)
        va.tick_labels.font.color.rgb = _TICK_RGB
    except Exception:
        pass

    # ── 7. Category axis: no gridlines, Trebuchet MS 10pt, #DDDDDD spine ─────
    try:
        ca = chart.category_axis
        ca.has_major_gridlines = False
        ca.has_minor_gridlines = False
        ca.format.line.color.rgb = _SPINE_RGB
        ca.tick_labels.font.name = "Trebuchet MS"
        ca.tick_labels.font.size = Pt(10)
        ca.tick_labels.font.color.rgb = _TICK_RGB
    except Exception:
        pass


def _fit_chart(chart_buf: io.BytesIO, max_w: float, max_h: float):
    """
    Return (width, height) as python-pptx EMU values sized to fit within
    max_w × max_h inches.  Exactly one of the two will be None so that
    python-pptx auto-scales the other dimension proportionally.
    Falls back to (Inches(max_w), None) on any error.
    """
    try:
        from PIL import Image as _PIL
        chart_buf.seek(0)
        with _PIL.open(chart_buf) as img:
            iw, ih = img.size
        if iw / ih >= max_w / max_h:
            return Inches(max_w), None     # width-constrained
        else:
            return None, Inches(max_h)     # height-constrained
    except Exception:
        return Inches(max_w), None


def generate_ppt(
    title: str,
    chart_buf: io.BytesIO,
    upload_date: str | None,
    *,
    chart_spec: dict | None = None,
    rows: list[dict] | None = None,
) -> tuple[io.BytesIO, str]:
    """
    Build a McKinsey-style 16:9 .pptx slide.
    When chart_spec + rows are supplied and the chart type is in _NATIVE_TYPES,
    inserts a native editable chart object.  Falls back to embedding the PNG
    image for unsupported types (funnel, heatmap, waterfall, treemap, etc.).
    Returns (BytesIO of the file, base filename without extension).
    """
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # White background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Title text
    txb = slide.shapes.add_textbox(Inches(0.4), Inches(0.22), Inches(12.5), Inches(0.52))
    tf  = txb.text_frame
    tf.word_wrap = False
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text           = title
    run.font.name      = "Trebuchet MS"
    run.font.bold      = True
    run.font.size      = Pt(18)
    run.font.color.rgb = _NAVY_RGB

    # Thin rule under title
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(0.4), Inches(0.80),
        Inches(12.93), Inches(0.80),
    )
    conn.line.color.rgb = _RULE_RGB
    conn.line.width     = Pt(0.8)

    # Chart area geometry
    chart_top  = 0.90
    chart_left = 0.40
    max_w = 13.33 - chart_left * 2   # 12.53"
    max_h = 7.50  - chart_top - 0.45  # leave 0.45" for footer

    left_emu   = Inches(chart_left)
    top_emu    = Inches(chart_top)
    width_emu  = Inches(max_w)
    height_emu = Inches(max_h)

    # ── Attempt native editable chart ────────────────────────────────────────
    native_ok = False
    if chart_spec and rows:
        chart_type_str = (chart_spec.get("chart_type") or "").lower()
        xl_type = _NATIVE_TYPES.get(chart_type_str)
        if xl_type is not None:
            try:
                if chart_type_str in ("scatter", "connected_scatter"):
                    cd = _make_xy_chart_data(chart_spec, rows, bubble=False)
                elif chart_type_str == "bubble":
                    cd = _make_xy_chart_data(chart_spec, rows, bubble=True)
                else:
                    cd = _make_category_chart_data(chart_spec, rows)

                if cd is not None:
                    gf = slide.shapes.add_chart(
                        xl_type, left_emu, top_emu, width_emu, height_emu, cd,
                    )
                    _style_chart(gf.chart, chart_type_str)
                    native_ok = True
                    logger.debug("Native pptx chart inserted: type=%s", chart_type_str)
            except Exception as exc:
                logger.warning(
                    "Native chart failed (%s), falling back to image: %s",
                    chart_type_str, exc,
                )

    # ── Fallback: embed matplotlib PNG ───────────────────────────────────────
    if not native_ok and chart_buf:
        pic_w, pic_h = _fit_chart(chart_buf, max_w, max_h)
        chart_buf.seek(0)
        slide.shapes.add_picture(
            chart_buf,
            left_emu, top_emu,
            width=pic_w, height=pic_h,
        )

    # Footer
    footer_txb = slide.shapes.add_textbox(
        Inches(0.4), Inches(7.12), Inches(12.5), Inches(0.3)
    )
    tf  = footer_txb.text_frame
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text           = _footer_text(upload_date)
    run.font.name      = "Trebuchet MS"
    run.font.size      = Pt(7)
    run.font.color.rgb = _GRAY_RGB

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf, _safe_filename(title)
